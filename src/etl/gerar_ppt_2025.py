
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Caminhos
DIRETORIO_2025 = r'C:\Users\franciscoj\Python_Initial\Pyhton_Web\data\base\2025'
ARQUIVO_EXCEL = os.path.join(DIRETORIO_2025, 'consolidado_2025_analise.xlsx')
ARQUIVO_PPT = os.path.join(DIRETORIO_2025, 'Apresentacao_Executiva_2025_v2.pptx')

def criar_slide_executivo():
    print("🚀 Iniciando geração da apresentação...")
    
    # Carregar dados
    try:
        df_resumo = pd.read_excel(ARQUIVO_EXCEL, sheet_name='Resumo_Executivo')
        df_setor = pd.read_excel(ARQUIVO_EXCEL, sheet_name='Resumo_Geral_Setor')
        df_picos = pd.read_excel(ARQUIVO_EXCEL, sheet_name='Top_Meses_Demanda')
    except Exception as e:
        print(f"❌ Erro ao ler Excel: {e}")
        return

    # Extrair valores chave
    total_chamados = df_resumo.loc[df_resumo['Métrica'] == 'Total Chamados (Ano)', 'Valor'].values[0]
    
    # Preparar PPT
    prs = Presentation()
    # Layout 16:9 (w=33.867 cm, h=19.05 cm)
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    # Slide 1: Dashboard Executivo
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # --- CABEÇALHO ---
    # Retângulo Azul Escuro (Fundo do Título)
    shape = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        Cm(0), Cm(0), Cm(33.87), Cm(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102) # Azul Corporativo Profundo
    shape.line.fill.background() # Sem borda
    
    # Título
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(20), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Panorama Executivo 2025: Análise de Demandas"
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtítulo (Data)
    txBoxSub = slide.shapes.add_textbox(Cm(25), Cm(0.8), Cm(8), Cm(1))
    tfSub = txBoxSub.text_frame
    pSub = tfSub.paragraphs[0]
    pSub.text = "Consolidado Anual"
    pSub.font.name = 'Arial'
    pSub.font.size = Pt(14)
    pSub.font.color.rgb = RGBColor(200, 200, 200)
    pSub.alignment = PP_ALIGN.RIGHT

    # --- KPI 1: VOLUME TOTAL (Esquerda) ---
    # Fundo Cinza Claro
    bg_kpi1 = slide.shapes.add_shape(1, Cm(1), Cm(3.5), Cm(9), Cm(6))
    bg_kpi1.fill.solid()
    bg_kpi1.fill.fore_color.rgb = RGBColor(240, 240, 240)
    bg_kpi1.line.color.rgb = RGBColor(200, 200, 200)
    
    # Título KPI
    box_title1 = slide.shapes.add_textbox(Cm(1.5), Cm(4), Cm(8), Cm(1))
    p1 = box_title1.text_frame.paragraphs[0]
    p1.text = "VOLUME TOTAL"
    p1.font.bold = True
    p1.font.size = Pt(14)
    p1.font.color.rgb = RGBColor(100, 100, 100)
    
    # Valor KPI
    box_val1 = slide.shapes.add_textbox(Cm(1.5), Cm(5), Cm(8), Cm(2))
    p_val1 = box_val1.text_frame.paragraphs[0]
    p_val1.text = str(total_chamados)
    p_val1.font.bold = True
    p_val1.font.size = Pt(54)
    p_val1.font.color.rgb = RGBColor(0, 51, 102) # Azul Forte
    p_val1.alignment = PP_ALIGN.CENTER
    
    # Subtexto KPI
    box_sub1 = slide.shapes.add_textbox(Cm(1.5), Cm(7.5), Cm(8), Cm(1))
    p_sub1 = box_sub1.text_frame.paragraphs[0]
    p_sub1.text = "Chamados Únicos\n(Deduplicados)"
    p_sub1.font.size = Pt(12)
    p_sub1.alignment = PP_ALIGN.CENTER
    p_sub1.font.color.rgb = RGBColor(80, 80, 80)

    # --- KPI 2: DISTRIBUIÇÃO (Centro) ---
    # Fundo Cinza Claro
    bg_kpi2 = slide.shapes.add_shape(1, Cm(11), Cm(3.5), Cm(10), Cm(6))
    bg_kpi2.fill.solid()
    bg_kpi2.fill.fore_color.rgb = RGBColor(240, 240, 240)
    bg_kpi2.line.color.rgb = RGBColor(200, 200, 200)

    # Título KPI
    box_title2 = slide.shapes.add_textbox(Cm(11.5), Cm(4), Cm(9), Cm(1))
    p2 = box_title2.text_frame.paragraphs[0]
    p2.text = "POR DEPARTAMENTO"
    p2.font.bold = True
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    
    # Tabela Simples Visual
    y_pos = 5.5
    for idx, row in df_setor.iterrows():
        setor = row['Setor']
        qtd = row['Qtd']
        perc = row['%'] * 100 # Assumindo que no Excel está decimal (0.76), se estiver texto '76.7%', precisa ajustar
        
        # Se for string (ex: '76.7%'), limpar
        if isinstance(perc, str):
            perc_str = perc
        else:
            perc_str = f"{perc:.1f}%"
            
        # Barra de Progresso Fundo
        bar_bg = slide.shapes.add_shape(1, Cm(11.5), Cm(y_pos+0.7), Cm(9), Cm(0.3))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(220, 220, 220)
        bar_bg.line.fill.background()
        
        # Barra de Progresso Valor
        # Converter perc para largura (max 9cm)
        try:
             perc_val = float(perc_str.replace('%', '')) / 100
        except:
             perc_val = 0.5
        
        bar_val = slide.shapes.add_shape(1, Cm(11.5), Cm(y_pos+0.7), Cm(9 * perc_val), Cm(0.3))
        bar_val.fill.solid()
        if 'ITI' in setor:
            bar_val.fill.fore_color.rgb = RGBColor(204, 0, 0) # Vermelho para ITI (Crítico)
        else:
            bar_val.fill.fore_color.rgb = RGBColor(0, 102, 204) # Azul para SPN
        bar_val.line.fill.background()
        
        # Texto
        box_item = slide.shapes.add_textbox(Cm(11.5), Cm(y_pos), Cm(9), Cm(0.8))
        p_item = box_item.text_frame.paragraphs[0]
        p_item.text = f"{setor}: {qtd} ({perc_str})"
        p_item.font.size = Pt(12)
        p_item.font.bold = True
        
        y_pos += 2.0

    # --- KPI 3: SAZONALIDADE (Direita) ---
    # Fundo Cinza Claro
    bg_kpi3 = slide.shapes.add_shape(1, Cm(22), Cm(3.5), Cm(10.8), Cm(6))
    bg_kpi3.fill.solid()
    bg_kpi3.fill.fore_color.rgb = RGBColor(240, 240, 240)
    bg_kpi3.line.color.rgb = RGBColor(200, 200, 200)

    # Título KPI
    box_title3 = slide.shapes.add_textbox(Cm(22.5), Cm(4), Cm(9), Cm(1))
    p3 = box_title3.text_frame.paragraphs[0]
    p3.text = "TOP 3 MESES (PICOS)"
    p3.font.bold = True
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(100, 100, 100)
    
    # Lista de Meses
    y_pos_mes = 5.5
    for idx, row in df_picos.iterrows():
        mes = row['Mes_Criacao']
        qtd = row['Qtd']
        
        # Traduzir Mês (opcional, simples)
        mapa_meses = {'August': 'Agosto', 'September': 'Setembro', 'November': 'Novembro', 'July': 'Julho'}
        mes_pt = mapa_meses.get(mes, mes)

        box_mes = slide.shapes.add_textbox(Cm(22.5), Cm(y_pos_mes), Cm(9), Cm(0.8))
        p_mes = box_mes.text_frame.paragraphs[0]
        p_mes.text = f"• {mes_pt}: {qtd} chamados"
        p_mes.font.size = Pt(14)
        if idx == 0: # Destaque para o primeiro
            p_mes.font.bold = True
            p_mes.font.color.rgb = RGBColor(204, 0, 0)
        
        y_pos_mes += 1.2

    # --- INSIGHT / RODAPÉ (Bottom) ---
    # Caixa de Destaque
    bg_insight = slide.shapes.add_shape(1, Cm(1), Cm(10.5), Cm(31.8), Cm(3))
    bg_insight.fill.solid()
    bg_insight.fill.fore_color.rgb = RGBColor(230, 240, 255) # Azul bem claro
    bg_insight.line.color.rgb = RGBColor(0, 51, 102)
    bg_insight.line.width = Pt(1.5)
    
    # Texto do Insight
    box_insight = slide.shapes.add_textbox(Cm(1.5), Cm(10.8), Cm(30), Cm(2.5))
    tf_insight = box_insight.text_frame
    p_ins = tf_insight.paragraphs[0]
    p_ins.text = "INSIGHT ESTRATÉGICO:"
    p_ins.font.bold = True
    p_ins.font.size = Pt(14)
    p_ins.font.color.rgb = RGBColor(0, 51, 102)
    
    p_ins_body = tf_insight.add_paragraph()
    p_ins_body.text = "A demanda apresentou forte concentração no 2º Semestre (Setembro/Novembro), indicando necessidade de planejamento preventivo para este período em 2026. O setor de Infraestrutura (ITI) absorve 3/4 de toda a carga de trabalho."
    p_ins_body.font.size = Pt(14)
    p_ins_body.space_before = Pt(6)

    # Rodapé Técnico
    txBoxFoot = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    pFoot = txBoxFoot.text_frame.paragraphs[0]
    pFoot.text = "Fonte: Base de Dados Backlog 2025 (Consolidado) | Francisco José Pereira"
    pFoot.font.size = Pt(9)
    pFoot.font.color.rgb = RGBColor(150, 150, 150)

    # Salvar
    prs.save(ARQUIVO_PPT)
    print(f"✅ Apresentação salva em: {ARQUIVO_PPT}")

if __name__ == "__main__":
    criar_slide_executivo()
