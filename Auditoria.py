import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import seaborn as sns

# Carregar os dados do Excel
def ler_excel(caminho_arquivo, aba_nome):
    df = pd.read_excel(caminho_arquivo, sheet_name=aba_nome)
    return df

# Função para limpar e converter valores para numéricos
def limpar_e_converter(df, colunas):
    for col in colunas:
        # Substituir espaços em branco ou células vazias por NaN
        df[col] = df[col].replace(r'^\s*$', pd.NA, regex=True)  # Substitui espaços em branco ou vazios por NaN
        df[col] = df[col].replace('%', '', regex=True).astype(float) / 100.0
    return df

# Gráfico de evolução de desempenho
def grafico_evolucao(df):
    plt.figure(figsize=(10, 6))
    paleta_cores = sns.color_palette("muted")  # Escolha de cores mais suaves e profissionais
    for categoria in df['Categoria'].unique():
        categoria_df = df[df['Categoria'] == categoria]
        plt.plot(categoria_df['ID'], categoria_df['2024'], label=f'{categoria} 2024', color=paleta_cores[0], linewidth=2)
        plt.plot(categoria_df['ID'], categoria_df['2023'], label=f'{categoria} 2023', color=paleta_cores[1], linewidth=2)
        plt.plot(categoria_df['ID'], categoria_df['2022'], label=f'{categoria} 2022', color=paleta_cores[2], linewidth=2)
        plt.plot(categoria_df['ID'], categoria_df['2021'], label=f'{categoria} 2021', color=paleta_cores[3], linewidth=2)
    
    plt.xlabel('ID', fontsize=14)
    plt.ylabel('Desempenho (%)', fontsize=14)
    plt.title('Evolução de Desempenho de Categorias ao Longo dos Anos', fontsize=16)
    plt.legend(title='Categorias', fontsize=12)
    plt.grid(True, linestyle='--', alpha=0.5)
    plt.tight_layout()
    plt.savefig('grafico_evolucao.png')
    plt.close()

# Gráfico de variações
def grafico_variacao(df):
    variacoes = ['Var. 2024-2023', 'Var. 2023-2022', 'Var. 2022-2021']
    # Limpar as colunas de variação e convertê-las
    df = limpar_e_converter(df, variacoes)
    
    df_variacoes = df[variacoes].mean()
    
    plt.figure(figsize=(8, 6))
    df_variacoes.plot(kind='bar', color='steelblue', edgecolor='black', linewidth=1.5)
    plt.xlabel('Variação Anual', fontsize=14)
    plt.ylabel('Valor Médio (%)', fontsize=14)
    plt.title('Média de Variação Anual entre os Anos', fontsize=16)
    plt.tight_layout()
    plt.savefig('grafico_variacao.png')
    plt.close()

# Criar apresentação PowerPoint
def criar_apresentacao():
    prs = Presentation()

    # Slide de título
    slide_titulo = prs.slides.add_slide(prs.slide_layouts[0])
    slide_titulo.shapes.title.text = "Relatório de Auditoria: Governança de TI"
    slide_titulo.placeholders[1].text = "Análise de Desempenho e Variação de Categorias de TI"

    # Slide de introdução
    slide_intro = prs.slides.add_slide(prs.slide_layouts[1])
    slide_intro.shapes.title.text = "Objetivos da Auditoria"
    slide_intro.shapes.placeholders[1].text = (
        "Esta auditoria visa avaliar o desempenho das categorias de TI nos últimos anos, identificando "
        "as principais variações e analisando a evolução dos indicadores de governança."
    )

    # Slide de gráfico de evolução
    slide_grafico_evolucao = prs.slides.add_slide(prs.slide_layouts[5])
    slide_grafico_evolucao.shapes.title.text = "Evolução de Desempenho"
    slide_grafico_evolucao.shapes.add_picture('grafico_evolucao.png', Inches(0.5), Inches(1.5), width=Inches(9))

    # Slide explicativo sobre gráfico de evolução
    slide_exp_evolucao = prs.slides.add_slide(prs.slide_layouts[1])
    slide_exp_evolucao.shapes.title.text = "Análise da Evolução"
    slide_exp_evolucao.shapes.placeholders[1].text = (
        "A análise de evolução mostra o desempenho das categorias ao longo dos anos. Podemos observar tendências "
        "de crescimento ou declínio, fundamentais para a definição de ações corretivas na governança de TI."
    )

    # Slide de gráfico de variações
    slide_grafico_variacao = prs.slides.add_slide(prs.slide_layouts[5])
    slide_grafico_variacao.shapes.title.text = "Variação Anual"
    slide_grafico_variacao.shapes.add_picture('grafico_variacao.png', Inches(0.5), Inches(1.5), width=Inches(9))

    # Slide explicativo sobre variações
    slide_exp_variacao = prs.slides.add_slide(prs.slide_layouts[1])
    slide_exp_variacao.shapes.title.text = "Análise das Variações"
    slide_exp_variacao.shapes.placeholders[1].text = (
        "A variação entre os anos indica a magnitude das mudanças no desempenho. Variações significativas precisam "
        "ser investigadas para entender os fatores que impactaram esses resultados."
    )

    # Slide de conclusões e recomendações
    slide_conclusao = prs.slides.add_slide(prs.slide_layouts[1])
    slide_conclusao.shapes.title.text = "Conclusões e Recomendações"
    slide_conclusao.shapes.placeholders[1].text = (
        "Com base na análise realizada, recomendamos a implementação de práticas de governança mais robustas "
        "para garantir a continuidade e melhoria no desempenho das categorias de TI. A gestão de mudanças e "
        "monitoramento contínuo das métricas são essenciais."
    )

    # Salvar apresentação
    prs.save('Relatorio_Auditoria_Governanca_TI.pptx')

# Caminho do arquivo Excel e nome da aba
caminho_arquivo_excel = "Base/planilha_auditoria.xlsx"
aba_nome = "Sheet"

# Ler dados do Excel
df = ler_excel(caminho_arquivo_excel, aba_nome)

# Gerar gráficos
graf
